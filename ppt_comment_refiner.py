import json
import requests
import argparse
from pptx import Presentation
from typing import List, Dict

class PPTCommentRefiner:
    """
    A class to refine comments in a PowerPoint presentation using OpenAI's GPT model.

    Attributes:
        ppt_path (str): Path to the input PowerPoint file.
        output_path (str): Path to save the refined PowerPoint file.
        api_key (str): OpenAI API key for authentication.
    """

    def __init__(self, ppt_path: str, output_path: str, api_key: str, base_url: str):
        self.ppt_path = ppt_path
        self.output_path = output_path
        self.api_key = api_key
        self.base_url = base_url
        self.model_name = "gpt-4o"
        self.headers = {
    "Content-Type": "application/json"
}

    def extract_notes(self) -> List[Dict[str, str]]:
        presentation = Presentation(self.ppt_path)
        notes_list = []
        for idx, slide in enumerate(presentation.slides):
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                content = slide.notes_slide.notes_text_frame.text.strip()
                if content:
                    notes_list.append({"index": idx, "content": content})
        return notes_list

    def refine_notes_with_gpt(self, notes_list: List[Dict[str, str]]) -> List[Dict[str, str]]:
        function = {
            "name": "refine_notes",
            "description": "Enhance presentation notes for better fluency and coherence.",
            "parameters": {
                "type": "object",
                "properties": {
                    "notes": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "index": {"type": "integer"},
                                "content": {"type": "string"}
                            },
                            "required": ["index", "content"]
                        }
                    }
                },
                "required": ["notes"]
            }
        }

        messages = [
            {"role": "system", "content": "You are a helpful assistant for improving presentation notes."},
            {"role": "user", "content": (
                "You will be provided with a list of notes extracted from a presentation. Each note is intended to be "
                "presented orally and should be enhanced for better fluency, correctness, coherence, and suitability for "
                "academic presentations. Please make the expressions natural, clear, suitable for speaking, and connect "
                "the ideas from different slides smoothly. Ensure that each refined note still matches the intent of the "
                "original but is more eloquent and fits a formal presentation style. The returned result must be in the "
                "same format as the provided input, with 'index' and 'content' keys."
            )}
        ]

        messages.append({"role": "user", "content": json.dumps({"notes": notes_list})})

        headers = self.headers.copy()
        headers["Authorization"] = f"Bearer {self.api_key}"

        data = {
            "model": self.model_name,
            "messages": messages,
            "functions": [function],
            "function_call": {"name": "refine_notes"}
        }

        try:
            response = requests.post(self.base_url, headers=headers, json=data)
            response.raise_for_status()
            result = response.json()
            function_args = result["choices"][0]["message"]["function_call"]["arguments"]
            refined_notes = json.loads(function_args)["notes"]
            return refined_notes
        except requests.exceptions.RequestException as e:
            print(f"An error occurred during the API request: {e}")
            return []

    def save_refined_notes_to_ppt(self, refined_notes: List[Dict[str, str]]):
        presentation = Presentation(self.ppt_path)
        for note in refined_notes:
            idx = note["index"]
            content = note["content"]
            slide = presentation.slides[idx]
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide.notes_slide.notes_text_frame.text = content
        presentation.save(self.output_path)

def main():
    parser = argparse.ArgumentParser(description="Refine PowerPoint presentation notes using OpenAI's GPT model.")
    parser.add_argument("--ppt_path", type=str, help="Path to the input PowerPoint file.")
    parser.add_argument("--output_path", type=str, help="Path to save the refined PowerPoint file.")
    parser.add_argument("--api_key", type=str, help="OpenAI API key for authentication.")
    parser.add_argument('--base_url', type=str, help='OpenAI API base URL')

    args = parser.parse_args()

    refiner = PPTCommentRefiner(args.ppt_path, args.output_path, args.api_key, args.base_url)

    # Extract notes from the presentation
    notes = refiner.extract_notes()

    # Refine notes using OpenAI's GPT model
    refined_notes = refiner.refine_notes_with_gpt(notes)

    # Save the refined notes back into the presentation
    refiner.save_refined_notes_to_ppt(refined_notes)

if __name__ == "__main__":
    main()


